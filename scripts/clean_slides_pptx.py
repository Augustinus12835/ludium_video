#!/usr/bin/env python3
"""
Clean PowerPoint lecture decks (.pptx) into content_cleaned.txt for the pipeline.

This is the slide-deck analogue of clean_book_chapter.py. Where that script handles
written prose (AsciiDoc/Markdown), this one handles *lecture slides* — terse bullet
fragments, exhibit tables, and accessibility ("ACCESS") spacing artifacts — and
rebuilds them into flowing, narration-friendly educational prose for a Technical-mode
video (Manim + SymPy downstream).

Publisher decks are usually copyrighted, so the cleaning prompt rewrites the material
substantially in its own words while preserving every formula, every worked-example
number, and the teaching order intact. Check the rights on your source deck before
distributing anything produced from it.

Workflow (the cleaning itself is authored by a Claude Code subagent — this script
makes no API call):
    inputs/slides/Course_Ch05.pptx  →  clean_slides_pptx.py --pipeline FIN --emit-prompt
        →  pipeline/FIN_Ch05_.../{slides_extracted.txt, clean_prompt.txt, source_info.json}
        →  subagent reads clean_prompt.txt → orchestrator writes content_cleaned.txt
        →  pipeline.py run ... --technical   (via /run-pipeline)

Usage:
    # Extract + scaffold the pipeline dir and render the cleaning prompt
    python scripts/clean_slides_pptx.py inputs/slides/Course_Ch05.pptx \
        --pipeline FIN --emit-prompt

    # Override chapter number / title
    python scripts/clean_slides_pptx.py inputs/slides/Course_Ch05.pptx --pipeline FIN \
        --emit-prompt --chapter 5 --title "Foreign Exchange Markets"
"""

import re
import sys
import json
import argparse
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))
from scripts.clean_book_chapter import chapter_title_to_dir_name


# ---------------------------------------------------------------------------
# Phase 1: PPTX extraction
# ---------------------------------------------------------------------------

# Lines that are pure slide chrome — dropped during extraction.
_CHROME_PATTERNS = [
    re.compile(r"^Slide(\s+\d+)?\s*$", re.IGNORECASE),   # "Slide 7" or bare "Slide"
    re.compile(r"^Access the text alternative", re.IGNORECASE),
    re.compile(r"^Return to parent slide", re.IGNORECASE),  # extended-desc nav chrome
    re.compile(r"©|All rights reserved", re.IGNORECASE),
    re.compile(r"^Source\s*:", re.IGNORECASE),       # exhibit data attribution
    re.compile(r"^Note\s*:", re.IGNORECASE),         # exhibit footnote
    re.compile(r"^\d{4}\s+Release\s*$"),             # "2026 Release" edition line
]


def _is_chrome(line: str) -> bool:
    s = line.strip()
    if not s:
        return True
    return any(p.search(s) for p in _CHROME_PATTERNS)


def _table_to_markdown(table) -> str:
    """Render a pptx table as a markdown table (raw cell text — the cleaning
    subagent/LLM normalizes ACCESS spacing, not a regex)."""
    rows = []
    for r in table.rows:
        cells = [c.text.strip().replace("\n", " ") for c in r.cells]
        rows.append(cells)
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    out = ["| " + " | ".join(rows[0]) + " |",
           "| " + " | ".join(["---"] * width) + " |"]
    for r in rows[1:]:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


def _clean_title(t: str) -> str:
    """Normalize a slide title: drop the in-frame vertical tab, strip a trailing
    continuation suffix the deck uses to spread one topic over several slides
    (e.g. 'Triangular Arbitrage Example 1' / 'Exhibit 5.2 ... 2' / '... 1(a)')."""
    t = t.replace("\x0b", " ").replace("\v", " ")
    t = re.sub(r"\s+", " ", t).strip()
    t = re.sub(r"\s+\d+(\([a-z]\))?$", "", t)   # "... 2", "... 1(a)"
    t = re.sub(r"\s+\([a-z]\)$", "", t)          # "... (b)"
    return t.strip()


def extract_pptx(path: Path, verbose: bool = False) -> tuple[str, str]:
    """Extract a deck into structured markdown-ish text.

    Returns (deck_title, body). Title placeholders become '## <title>' headings;
    bullet text and tables follow. Consecutive slides that share a (normalized)
    title are merged under one heading so split topics rejoin. Slide chrome
    (page numbers, copyright, accessibility lines, exhibit source/notes) is dropped.
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    deck_title = None
    blocks: list[str] = []
    last_heading = None

    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        _SLIDE_NUM = PP_PLACEHOLDER.SLIDE_NUMBER
    except Exception:
        _SLIDE_NUM = None

    for idx, slide in enumerate(prs.slides, 1):
        title_shape = slide.shapes.title
        # python-pptx hands back a fresh wrapper per access, so identity (`is`)
        # fails — match the title by its stable shape_id instead.
        title_id = title_shape.shape_id if title_shape else None
        raw_title = title_shape.text if title_shape else ""
        heading = _clean_title(raw_title) if raw_title.strip() else None

        if idx == 1:
            # Title card: "Chapter Five\x0bThe Market for Foreign Exchange".
            t = raw_title.replace("\x0b", "\n").replace("\v", "\n")
            parts = [p.strip() for p in t.split("\n") if p.strip()]
            parts = [p for p in parts
                     if not re.match(r"^Chapter\s+[A-Za-z0-9]+$", p, re.IGNORECASE)]
            deck_title = " ".join(parts).strip() or None
            continue  # title card carries no teaching content

        body_lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False) and shape.has_table:
                md = _table_to_markdown(shape.table)
                if md:
                    body_lines.append(md)
                continue
            if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
                continue
            if title_id is not None and shape.shape_id == title_id:
                continue
            if (_SLIDE_NUM is not None and shape.is_placeholder
                    and shape.placeholder_format.type == _SLIDE_NUM):
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if not line:
                    line = para.text.strip()
                if line and not _is_chrome(line):
                    body_lines.append(line)

        # Drop a body line that merely echoes the slide title.
        if heading and body_lines and _clean_title(body_lines[0]).lower() == heading.lower():
            body_lines = body_lines[1:]

        body = "\n".join(body_lines).strip()
        if not heading and not body:
            continue

        if heading and heading == last_heading:
            # continuation slide — append body under the same heading
            if body:
                blocks.append(body)
        else:
            if heading:
                blocks.append(f"\n## {_clean_title(heading)}")
                last_heading = heading
            if body:
                blocks.append(body)

    text = "\n".join(blocks).strip()
    text = re.sub(r"\n{3,}", "\n\n", text)
    if verbose:
        print(f"      extracted {len(prs.slides._sldIdLst)} slides → {len(text):,} chars")
    return deck_title, text


# ---------------------------------------------------------------------------
# ACCESS spacing artifacts (the "_ACCESS" decks space out acronyms/codes/years
# so screen readers spell them: "F X", "U S D", "20 25") are NOT normalized with
# regex here — a deterministic pass is brittle ("C N Y" vs "N Y", title leaks).
# Extraction stays raw and the cleaning AGENT normalizes them; its prompt lists
# the exact patterns. This keeps the deterministic layer to only what must be
# deterministic: pulling text/tables out of the binary .pptx.
# ---------------------------------------------------------------------------


# ---------------------------------------------------------------------------
# Phase 2: cleaning prompt — finance lecture-slides profile
# Rendered via --emit-prompt to a Claude Code subagent, which returns the cleaned
# content for the orchestrator to write to content_cleaned.txt.
# ---------------------------------------------------------------------------

FINANCE_SYSTEM_PROMPT = """\
You are an educational content editor preparing a finance textbook's lecture slides for \
an automated video pipeline that renders equations and worked examples with Manim and \
verifies the arithmetic with SymPy. The source is a publisher's copyrighted lecture deck. \
Rewrite it substantially in your own words — never \
mirror the slide phrasing — while preserving every formula, every number, and the teaching \
order with complete accuracy."""

FINANCE_CLEANING_PROMPT = """\
Rewrite this set of lecture slides into clean, narration-friendly educational prose for a \
video script pipeline. A narrator will speak it; equations, tables, and worked examples are \
animated. The slides are terse bullet fragments spread across many slides — your job is to \
rebuild them into flowing, coherent teaching prose. Accuracy of every formula and number is \
paramount.

## FIX accessibility ("ACCESS") spacing artifacts:
The deck spells out acronyms, currency codes, and years with spaces so screen readers read \
them letter-by-letter. Restore the normal forms in your prose:
- "F X" -> "FX",  "O T C" -> "OTC",  "U. S." / "U S" -> "U.S.",  "N Y" -> "NY"
- Currency codes: "U S D" -> "USD", "G B P" -> "GBP", "J P Y" -> "JPY", "C N Y" -> "CNY", \
"C H F" -> "CHF", "S Fr" / "S F r" -> "SFr"
- "U B S" -> "UBS", "J P Morgan" -> "JPMorgan", "B I S" -> "BIS"
- Split years: "20 25" -> "2025"
Write currency amounts normally ("USD 10 million", "$1,000,000").

## REMOVE completely:
- Slide chrome that survived extraction: "Slide N", "Access the text alternative...", \
publisher copyright lines, author names, "Source:"/"Note:" exhibit footnotes
- Continuation markers in headings ("Example 1", "1(a)", "... 2") — MERGE the split slides \
into one continuous treatment of that topic
- Pure cross-references to other chapters ("as we saw in Chapter 3", "see Exhibit 7.4")

## KEEP and preserve EXACTLY (do not drop, do not simplify):
- **Every formula.** Keep it as LaTeX — inline as $...$, displayed as $$...$$. Use real \
notation for sub/superscripts and quote currencies, e.g. spot/cross/forward rates \
$S(\\$/\\pounds)$, $S(j/k) = S(\\$/k)\\times S(j/\\$)$, bid/ask $S^b$, $S^a$, forward \
$F_N(\\$/\\text{{SFr}})$, the annualized forward premium formula. Do NOT paraphrase an \
equation into words — keep the symbolic form.
- **Every worked example in full**, with ALL of its numbers and every intermediate step \
(Example 5.2 cross-rate bid-ask spread; the cross-rate FX transactions; the triangular \
arbitrage example with each leg and the final profit). These become the step-by-step math \
frames — losing a step or a number breaks them. Keep each computed value; if a number on a \
slide looks internally inconsistent, keep the slide's value but you may note the apparent \
discrepancy in a parenthetical so the reviewer can check it.
- All definitions, market conventions (big figure / small figure, "12 to 17" quoting, \
mid-rate), and the logical teaching order.
- The pedagogically essential figures from exhibit tables (e.g. FX is the world's largest \
market at roughly USD 9.6 trillion daily turnover; the USD is on one side of ~88% of all \
trades). Convert a data table into one or two sentences of prose highlighting the key \
numbers and what they show — do not reproduce the whole table verbatim, but keep the \
headline figures accurate.

## TRANSFORM:
- Bullet fragments -> flowing prose at a narration-friendly pace (shorter sentences), \
WITHOUT removing teaching content
- A table the example actually computes from (bid/ask quote grids) -> keep the specific \
quotes inline in the prose so the worked example is self-contained
- Passive -> active voice where natural; use your own sentence structure

## REWRITE guidelines:
- Organize as flowing educational prose under clear topic headings (keep the deck's section \
structure: function & structure of the FX market, market participants, the spot market, \
quotations and cross-rates, the bid-ask spread and spot trading, triangular arbitrage, \
market microstructure, the forward market and forward premium).
- Be comprehensive — this is the single source of truth for the whole video.
- Preserve technical accuracy completely.

---

LECTURE SLIDES:
{text}

---

CLEANED EDUCATIONAL CONTENT:"""


def build_clean_prompt(raw: str) -> tuple[str, str]:
    """Return (system, user) — the exact cleaning prompt for the cleaning subagent."""
    return FINANCE_SYSTEM_PROMPT, FINANCE_CLEANING_PROMPT.format(text=raw)


# ---------------------------------------------------------------------------
# Phase 3: Output
# ---------------------------------------------------------------------------

def chapter_num_from_filename(path: Path) -> int | None:
    m = re.search(r"[Cc]h(?:apter)?[_\-]?(\d{1,2})", path.stem)
    return int(m.group(1)) if m else None


def build_source_info(textbook: str, num: int, title: str | None) -> dict:
    return {
        "type": "textbook_slides",
        "textbook": textbook,
        "chapter": num,
        "chapter_title": title,
        "note": ("Derived and substantially rewritten from publisher lecture slides. "
                 "If the source deck is copyrighted, check your rights before "
                 "distributing anything produced from it."),
    }


DEFAULT_TEXTBOOK = "Textbook lecture slides"


def emit_prompt_for(deck: Path, args) -> None:
    """Scaffold the pipeline dir and render the cleaning prompt for a subagent.

    Writes slides_extracted.txt, source_info.json, and clean_prompt.txt (the full
    system+user prompt) into pipeline/<dir>/, then prints a COMPACT meta JSON. The
    orchestrator spawns one subagent told to read clean_prompt.txt and return
    only the cleaned content, then writes that (with the header) to content_cleaned.txt.
    """
    deck_title, raw = extract_pptx(deck, args.verbose)
    if not raw.strip():
        print(json.dumps({"error": f"no extractable text in {deck}"}))
        sys.exit(1)
    num = args.chapter or chapter_num_from_filename(deck) or 0
    title = args.title or deck_title
    system, user = build_clean_prompt(raw)

    if not args.pipeline:
        print("Error: --emit-prompt requires --pipeline PREFIX")
        sys.exit(1)
    dir_name = chapter_title_to_dir_name(args.pipeline, num, title)
    pdir = Path("pipeline") / dir_name
    pdir.mkdir(parents=True, exist_ok=True)
    (pdir / "slides_extracted.txt").write_text(raw, encoding="utf-8")
    (pdir / "source_info.json").write_text(
        json.dumps(build_source_info(args.textbook, num, title), indent=2, ensure_ascii=False),
        encoding="utf-8")
    (pdir / "clean_prompt.txt").write_text(f"{system}\n\n{user}", encoding="utf-8")

    header = f"<!-- Source: {args.textbook}, Chapter {num} -->"
    meta = {
        "pipeline_dir": str(pdir),
        "content_path": str(pdir / "content_cleaned.txt"),
        "clean_prompt": str(pdir / "clean_prompt.txt"),
        "slides_extracted": str(pdir / "slides_extracted.txt"),
        "chapter": num,
        "title": title,
        "content_header": header,
        "content_title_line": f"# {title}",
        "extracted_chars": len(raw),
        "chunking_needed": len(raw) > 45000,
    }
    print(json.dumps(meta, ensure_ascii=False, indent=2))


def main():
    ap = argparse.ArgumentParser(
        description="Extract PowerPoint lecture decks and render the cleaning prompt "
                    "for a Claude Code subagent (content_cleaned.txt is subagent-authored; "
                    "no API call is made)")
    ap.add_argument("input", help="A .pptx deck or a directory of decks")
    ap.add_argument("--pipeline", metavar="PREFIX",
                    help="Create pipeline/<PREFIX>_Ch<NN>_<Title>/ (required with --emit-prompt)")
    ap.add_argument("--chapter", type=int, help="Override chapter number")
    ap.add_argument("--title", help="Override chapter title (only valid for a single deck)")
    ap.add_argument("--textbook", default=DEFAULT_TEXTBOOK,
                    help="Textbook citation for source_info.json, e.g. "
                         "'Author, Title (Edition), Publisher' (default: %(default)r)")
    ap.add_argument("--emit-prompt", action="store_true",
                    help="Extract + scaffold the pipeline dir, then print the {system,user,meta} "
                         "cleaning prompt as JSON for a Claude Code subagent (no API call). The "
                         "subagent's returned text is written to content_cleaned.txt by the "
                         "orchestrator.")
    ap.add_argument("--verbose", action="store_true")
    args = ap.parse_args()

    input_path = Path(args.input)
    if input_path.is_file():
        decks = [input_path]
    elif input_path.is_dir():
        decks = sorted(f for f in input_path.iterdir() if f.suffix.lower() == ".pptx")
    else:
        print(f"Error: {input_path} not found")
        sys.exit(1)

    if not decks:
        print(f"No .pptx files found in {input_path}")
        sys.exit(1)
    if args.title and len(decks) > 1:
        print("Error: --title only valid for a single deck")
        sys.exit(1)

    if not args.emit_prompt:
        sys.exit(
            "Deck cleaning is authored by a Claude Code subagent (no API call is made).\n"
            "Run with --emit-prompt (one deck at a time): python scripts/clean_slides_pptx.py "
            f"{decks[0]} --pipeline <PREFIX> --emit-prompt\n"
            "then spawn a subagent on the printed clean_prompt.txt and write its output to "
            "content_cleaned.txt — see .claude/skills/run-pipeline/references/sources.md."
        )

    if len(decks) != 1:
        print("Error: --emit-prompt handles one deck at a time")
        sys.exit(1)
    emit_prompt_for(decks[0], args)


if __name__ == "__main__":
    main()
